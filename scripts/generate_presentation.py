"""
Generate a professional PowerPoint presentation for the
AI Agent Ticket Triage & Auto-Routing System.

Install dependencies first:
    pip install python-pptx comtypes

Run:
    python scripts/generate_presentation.py

Output:
    AI_Ticket_Triage_Presentation.pptx  — modern format (opens in any version)
    AI_Ticket_Triage_Presentation.ppt   — legacy .ppt format (via PowerPoint COM)

Note: .ppt conversion requires Microsoft PowerPoint to be installed on the machine.
      If PowerPoint is not installed, the .pptx file is produced and can be manually
      saved as .ppt via File → Save As in PowerPoint.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Color Palette ─────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x00, 0x15, 0x29)   # sidebar dark
BLUE       = RGBColor(0x16, 0x77, 0xFF)   # AntD primary blue
LIGHT_BLUE = RGBColor(0xE6, 0xF4, 0xFF)   # light blue bg
GREEN      = RGBColor(0x52, 0xC4, 0x1A)
ORANGE     = RGBColor(0xFA, 0x8C, 0x16)
RED        = RGBColor(0xF5, 0x22, 0x2D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x26, 0x26, 0x26)
GREY_TEXT  = RGBColor(0x59, 0x59, 0x59)
LIGHT_GREY = RGBColor(0xF0, 0xF2, 0xF5)

# ── Slide Dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


# ── Helper: fill slide background ─────────────────────────────────────────────
def fill_bg(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Helper: add rectangle shape ───────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill: RGBColor,
             text="", font_size=18, font_color=WHITE, bold=False,
             line_color=None, line_width=None, align=PP_ALIGN.CENTER,
             font_name="Calibri"):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = font_name
    return shape


# ── Helper: add text box ───────────────────────────────────────────────────────
def add_text(slide, left, top, width, height, text,
             size=18, color=DARK_TEXT, bold=False, italic=False,
             align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox


def add_bullet_text(slide, left, top, width, height,
                    title, bullets: list[tuple[str, str]],
                    title_size=20, bullet_size=16,
                    title_color=NAVY, bullet_color=DARK_TEXT):
    """Add a text box with a bold title then bullet points."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title line
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(title_size)
    r0.font.color.rgb = title_color
    r0.font.bold = True
    r0.font.name = "Calibri"

    for icon, line in bullets:
        para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(4)
        run = para.add_run()
        run.text = f"  {icon}  {line}"
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = bullet_color
        run.font.name = "Calibri"
    return txBox


# ── Helper: add arrow ────────────────────────────────────────────────────────
def add_arrow_down(slide, cx, top, height, color=BLUE):
    """Add a downward arrow connector."""
    from pptx.util import Pt
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        cx, top, cx, top + height
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: Title Slide
# ─────────────────────────────────────────────────────────────────────────────
def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    # Left accent bar
    add_rect(sl, Inches(0), Inches(0), Inches(0.4), H, BLUE)

    # Main title
    add_text(sl, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
             "AI Agent Ticket Triage &\nAuto-Routing System",
             size=44, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Subtitle
    add_text(sl, Inches(0.8), Inches(3.6), Inches(10), Inches(0.8),
             "Automated ITSM Intelligence — Classify · Prioritize · Route",
             size=22, color=RGBColor(0x8C, 0xB8, 0xFF), align=PP_ALIGN.LEFT)

    # Tag line
    add_text(sl, Inches(0.8), Inches(4.6), Inches(10), Inches(0.6),
             "Built with  Python · FastAPI · React · Ollama · scikit-learn",
             size=14, color=RGBColor(0x59, 0x6B, 0x8C), align=PP_ALIGN.LEFT)

    # Date
    add_text(sl, Inches(0.8), Inches(6.6), Inches(5), Inches(0.5),
             "June 2026",
             size=13, color=RGBColor(0x59, 0x6B, 0x8C), align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: Problem Statement
# ─────────────────────────────────────────────────────────────────────────────
def slide_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl, LIGHT_GREY)

    # Header bar
    add_rect(sl, Inches(0), Inches(0), W, Inches(1.2), NAVY)
    add_text(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             "The Problem", size=30, color=WHITE, bold=True)

    problems = [
        ("⏱", "Manual triage takes 15–30 minutes per ticket batch"),
        ("❌", "Inconsistent routing — wrong team assigned, rework required"),
        ("📉", "No priority intelligence — critical issues missed in the queue"),
        ("🔁", "Repetitive work — analysts spend 40%+ of time on routing"),
        ("📊", "No visibility into ticket trends or team workload"),
    ]

    for i, (icon, text) in enumerate(problems):
        y = Inches(1.5) + i * Inches(0.95)
        add_rect(sl, Inches(0.5), y, Inches(12), Inches(0.75),
                 WHITE, line_color=RGBColor(0xD9, 0xD9, 0xD9), line_width=Pt(1))
        add_text(sl, Inches(0.8), y + Inches(0.12), Inches(11.5), Inches(0.6),
                 f"{icon}  {text}", size=18, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: Our Solution
# ─────────────────────────────────────────────────────────────────────────────
def slide_solution(prs):
    sl = blank_slide(prs)
    fill_bg(sl, LIGHT_GREY)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.2), BLUE)
    add_text(sl, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             "Our Solution — AI Agent Pipeline", size=30, color=WHITE, bold=True)

    add_text(sl, Inches(0.5), Inches(1.4), Inches(12.5), Inches(0.7),
             "An intelligent multi-agent system that automatically classifies, prioritizes,\n"
             "and routes IT support tickets — with zero manual intervention for high-confidence cases.",
             size=16, color=GREY_TEXT)

    boxes = [
        (BLUE,   "🔍  Classify",     "Determines ticket\ncategory automatically"),
        (ORANGE, "⚡  Prioritize",   "Assigns P1–P4 priority\nbased on impact"),
        (GREEN,  "📬  Route",        "Assigns to the right\nteam & person"),
        (NAVY,   "📊  Confidence",   "Decides: auto-route,\nflag, or hold"),
        (RGBColor(0x72, 0x2E, 0xD1), "🔄  Learn", "Feedback loop\nimproves accuracy"),
    ]

    for i, (color, title, desc) in enumerate(boxes):
        x = Inches(0.4) + i * Inches(2.55)
        add_rect(sl, x, Inches(2.5), Inches(2.4), Inches(1.0),
                 color, text=title, font_size=16, bold=True)
        add_text(sl, x, Inches(3.65), Inches(2.4), Inches(0.8),
                 desc, size=13, color=GREY_TEXT, align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(sl, x + Inches(2.4), Inches(2.9), Inches(0.15), Inches(0.5),
                     "→", size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Bottom outcome
    add_rect(sl, Inches(1), Inches(5.0), Inches(11), Inches(1.2),
             RGBColor(0xF6, 0xFF, 0xED), line_color=GREEN, line_width=Pt(2))
    add_text(sl, Inches(1.2), Inches(5.15), Inches(10.6), Inches(0.9),
             "✅  Result: >85% of tickets routed automatically — "
             "reducing manual effort by ≥50% and eliminating routing errors",
             size=16, color=RGBColor(0x13, 0x52, 0x00), bold=False, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: End-to-End Workflow
# ─────────────────────────────────────────────────────────────────────────────
def slide_workflow(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.1), NAVY)
    add_text(sl, Inches(0.5), Inches(0.18), Inches(12), Inches(0.75),
             "End-to-End Workflow — How It Works", size=28, color=WHITE, bold=True)

    # Row 1: Input
    add_rect(sl, Inches(0.3), Inches(1.3), Inches(2.5), Inches(1.0),
             LIGHT_BLUE, text="📄  Excel File\n(tickets.xlsx)", font_size=14,
             font_color=NAVY, bold=True, line_color=BLUE, line_width=Pt(1.5))

    add_text(sl, Inches(2.85), Inches(1.65), Inches(0.8), Inches(0.4),
             "Upload", size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(2.85), Inches(1.8), Inches(0.8), Inches(0.4),
             "━━━▶", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(3.7), Inches(1.3), Inches(2.8), Inches(1.0),
             BLUE, text="🌐  Web Portal\n(React + Ant Design)", font_size=13,
             font_color=WHITE, bold=False)

    add_text(sl, Inches(6.55), Inches(1.65), Inches(0.9), Inches(0.4),
             "Submit", size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(6.55), Inches(1.8), Inches(0.9), Inches(0.4),
             "━━━▶", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(7.5), Inches(1.3), Inches(2.8), Inches(1.0),
             NAVY, text="⚙️  FastAPI\nBackend", font_size=13,
             font_color=WHITE, bold=False)

    add_text(sl, Inches(10.35), Inches(1.65), Inches(0.7), Inches(0.4),
             "Triggers", size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(10.35), Inches(1.8), Inches(0.7), Inches(0.4),
             "━━━▶", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(11.1), Inches(1.3), Inches(1.9), Inches(1.0),
             RGBColor(0x72, 0x2E, 0xD1), text="🤖  AI\nAgents", font_size=13,
             font_color=WHITE, bold=False)

    # Arrow down to agents box
    add_text(sl, Inches(11.7), Inches(2.35), Inches(0.5), Inches(0.5),
             "▼", size=16, color=RGBColor(0x72, 0x2E, 0xD1), bold=True, align=PP_ALIGN.CENTER)

    # Row 2: Agent pipeline
    agents = [
        ("1\n🔍\nClassify", BLUE),
        ("2\n⚡\nPriority", ORANGE),
        ("3\n📬\nRoute", GREEN),
        ("4\n📊\nConfidence", NAVY),
    ]
    for i, (label, color) in enumerate(agents):
        x = Inches(3.0) + i * Inches(2.45)
        add_rect(sl, x, Inches(2.9), Inches(2.2), Inches(1.1),
                 color, text=label, font_size=13, font_color=WHITE, bold=True)
        if i < 3:
            add_text(sl, x + Inches(2.2), Inches(3.3), Inches(0.25), Inches(0.5),
                     "▶", size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_text(sl, Inches(0.5), Inches(3.1), Inches(2.4), Inches(0.9),
             "ML Model\n(scikit-learn)\n+  Ollama LLM", size=12,
             color=GREY_TEXT, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(2.6), Inches(3.35), Inches(0.5), Inches(0.4),
             "━━▶", size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Arrow down to decision
    add_text(sl, Inches(11.2), Inches(4.05), Inches(1.5), Inches(0.4),
             "▼  Decision", size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Row 3: Decision outcomes
    outcomes = [
        (Inches(1.5),  GREEN,  "✅  AUTO-ROUTED\n> 85% confidence\n\nEmail sent automatically"),
        (Inches(5.2),  ORANGE, "⚠️  FLAGGED\n65–85% confidence\n\nNeeds human review"),
        (Inches(8.9),  RED,    "🔴  HELD\n< 65% confidence\n\nManual classification"),
    ]
    for x, color, text in outcomes:
        add_rect(sl, x, Inches(4.6), Inches(3.3), Inches(1.3),
                 color, text=text, font_size=13, font_color=WHITE, bold=False)

    # Enriched Excel output note
    add_rect(sl, Inches(0.3), Inches(6.1), Inches(12.7), Inches(0.7),
             LIGHT_BLUE, line_color=BLUE, line_width=Pt(1),
             text="📥  Enriched Excel downloaded  |  AI_Category · AI_Priority · Assigned_Team · Assigned_To · Confidence_Score · Routing_Status · Email_Sent",
             font_size=13, font_color=NAVY, bold=False)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: Technology Stack
# ─────────────────────────────────────────────────────────────────────────────
def slide_tech(prs):
    sl = blank_slide(prs)
    fill_bg(sl, LIGHT_GREY)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.1), NAVY)
    add_text(sl, Inches(0.5), Inches(0.18), Inches(12), Inches(0.75),
             "Technology Stack — 100% Open Source", size=28, color=WHITE, bold=True)

    categories = [
        ("🤖  AI / ML",     BLUE,   ["Ollama (gpt-oss:120b-cloud) — Local LLM inference",
                                      "scikit-learn — TF-IDF + Logistic Regression classifier",
                                      "sentence-transformers — Text embeddings (future)"]),
        ("⚙️  Backend",     NAVY,   ["Python 3.13 + FastAPI — REST API",
                                      "SQLAlchemy + aiosqlite — Async database layer",
                                      "Jinja2 — HTML email templates"]),
        ("🎨  Frontend",    RGBColor(0x72, 0x2E, 0xD1),
                                     ["React 18 + Vite — Modern SPA",
                                      "Ant Design 5 — Enterprise UI components",
                                      "Recharts — KPI charts and visualisations"]),
        ("📧  Cloud Action", GREEN,  ["Gmail SMTP (smtplib) — Real email on auto-route",
                                      "App Password authentication — Secure, no OAuth needed",
                                      "Jinja2 HTML templates — Rich formatted emails"]),
        ("📊  Monitoring",  ORANGE,  ["Prometheus — Metrics collection",
                                      "Grafana — Dashboard visualisation",
                                      "Structured Python logging — Full audit trail"]),
        ("☁️  Deployment",  RGBColor(0x08, 0x97, 0x9C),
                                     ["Azure Container Apps — Serverless containers",
                                      "Terraform — Infrastructure as Code",
                                      "GitHub Actions — CI/CD pipeline"]),
    ]

    cols = 3
    for i, (cat, color, items) in enumerate(categories):
        row = i // cols
        col = i % cols
        x = Inches(0.3) + col * Inches(4.35)
        y = Inches(1.3) + row * Inches(2.6)
        add_rect(sl, x, y, Inches(4.1), Inches(0.5), color, text=cat,
                 font_size=15, bold=True)
        for j, item in enumerate(items):
            add_text(sl, x + Inches(0.1), y + Inches(0.55) + j * Inches(0.55),
                     Inches(4.0), Inches(0.5),
                     f"•  {item}", size=12, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: Multi-Agent Pipeline Detail
# ─────────────────────────────────────────────────────────────────────────────
def slide_agents(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.1), BLUE)
    add_text(sl, Inches(0.5), Inches(0.18), Inches(12), Inches(0.75),
             "Multi-Agent Pipeline — Intelligence Inside", size=28, color=WHITE, bold=True)

    agents = [
        (BLUE,   "1.  Classifier Agent",
                 "What type of problem is this?",
                 "• ML model (TF-IDF + Logistic Regression)\n• Ollama LLM fallback for ambiguous tickets\n• Outputs: category + probability score"),
        (ORANGE, "2.  Priority Agent",
                 "How urgent is it?",
                 "• Keyword rules for obvious cases (P1: 'outage', 'down')\n• Ollama reasoning for ambiguous tickets\n• Outputs: P1 Critical / P2 High / P3 Medium / P4 Low"),
        (GREEN,  "3.  Routing Agent",
                 "Who should handle it?",
                 "• Reads routing_config.xlsx (editable without code)\n• P1/P2 → Team Lead automatically\n• P3/P4 → Ollama picks best-fit member"),
        (NAVY,   "4.  Confidence Agent",
                 "Are we sure enough to act?",
                 "• Weighted score: ML (50%) + LLM (30%) + Priority (20%)\n• >85% → Auto-Route + Email\n• 65-85% → Flagged for review\n• <65% → Held for manual"),
        (RGBColor(0x72, 0x2E, 0xD1), "5.  Feedback Agent",
                 "How do we get smarter?",
                 "• Stores human corrections in SQLite\n• Triggers retraining when N corrections reached\n• Continuously improves accuracy over time"),
    ]

    for i, (color, title, subtitle, detail) in enumerate(agents):
        x = Inches(0.2) + i * Inches(2.62)
        add_rect(sl, x, Inches(1.2), Inches(2.5), Inches(0.65),
                 color, text=title, font_size=13, bold=True)
        add_text(sl, x, Inches(1.9), Inches(2.5), Inches(0.5),
                 subtitle, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_rect(sl, x, Inches(2.45), Inches(2.5), Inches(4.6),
                 LIGHT_GREY, line_color=RGBColor(0xD9, 0xD9, 0xD9),
                 line_width=Pt(1))
        add_text(sl, x + Inches(0.1), Inches(2.6), Inches(2.3), Inches(4.3),
                 detail, size=12, color=DARK_TEXT, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: Supported Teams
# ─────────────────────────────────────────────────────────────────────────────
def slide_teams(prs):
    sl = blank_slide(prs)
    fill_bg(sl, LIGHT_GREY)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.1), NAVY)
    add_text(sl, Inches(0.5), Inches(0.18), Inches(12), Inches(0.75),
             "Supported Teams & Routing Configuration", size=28, color=WHITE, bold=True)

    teams = [
        (BLUE,   "🧪  Testing Team",    "Test failures, UAT issues,\nautomation errors, QA sign-offs"),
        (GREEN,  "📡  EDI Team",         "EDI transactions, 835/837 files,\nHIPAA validation, partner setup"),
        (ORANGE, "⚙️  BizTalk Team",     "Orchestration failures, pipeline\nexceptions, adapter config"),
        (RED,    "🗄️  Database Team",    "Connection timeouts, slow queries,\ndeadlocks, backup failures"),
        (RGBColor(0x08, 0x97, 0x9C), "🌐  Network Team", "VPN, DNS, firewall rules,\nload balancer issues"),
        (RGBColor(0x72, 0x2E, 0xD1), "🔐  Access Team",  "Account lockouts, password\nresets, permission issues"),
        (RGBColor(0xEB, 0x2F, 0x96), "🛡️  Security Team", "Phishing, malware, compliance\naudits, patch deployment"),
    ]

    for i, (color, name, desc) in enumerate(teams):
        row = i // 4
        col = i % 4
        x = Inches(0.4) + col * Inches(3.25)
        y = Inches(1.4) + row * Inches(2.2)
        add_rect(sl, x, y, Inches(3.0), Inches(0.55),
                 color, text=name, font_size=14, bold=True)
        add_rect(sl, x, y + Inches(0.55), Inches(3.0), Inches(1.35),
                 WHITE, text=desc, font_size=13, font_color=GREY_TEXT,
                 line_color=RGBColor(0xD9, 0xD9, 0xD9), line_width=Pt(1))

    add_rect(sl, Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.8),
             LIGHT_BLUE, line_color=BLUE, line_width=Pt(1.5),
             text="🔧  Need to add a new team? Just update data/routing_config.xlsx — no code changes required",
             font_size=14, font_color=NAVY, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8: KRA Achievements
# ─────────────────────────────────────────────────────────────────────────────
def slide_kra(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.1), GREEN)
    add_text(sl, Inches(0.5), Inches(0.18), Inches(12), Inches(0.75),
             "KRA Achievements", size=28, color=WHITE, bold=True)

    kras = [
        (GREEN,  "✅", "Data Collection & Baseline",
                 "50 synthetic labeled tickets across 7 teams generated.\n"
                 "Baseline metrics established before model training."),
        (GREEN,  "✅", "Classifier Model — 70–80% Accuracy",
                 "TF-IDF + Logistic Regression pipeline trained.\n"
                 "Accuracy, F1-macro and confusion matrix reported."),
        (GREEN,  "✅", "Multi-Agent Architecture",
                 "5 specialised agents: Classifier, Priority, Routing,\n"
                 "Confidence, Feedback — all connected in one pipeline."),
        (GREEN,  "✅", "Excel Auto-Routing Pipeline",
                 "Upload .xlsx → enrich with AI columns → download.\n"
                 "Confidence thresholds: >85% auto-route, else flag/hold."),
        (GREEN,  "✅", "Real Cloud Action — Gmail SMTP",
                 "Auto-routed tickets trigger email notification\n"
                 "to team assignee. Live email delivered."),
        (ORANGE, "🔄", "Feedback Loop & Similarity (In Progress)",
                 "Human corrections stored in SQLite.\n"
                 "Retraining triggered at configurable threshold."),
    ]

    cols = 2
    for i, (color, icon, title, detail) in enumerate(kras):
        row = i // cols
        col = i % cols
        x = Inches(0.4) + col * Inches(6.5)
        y = Inches(1.3) + row * Inches(1.8)
        add_rect(sl, x, y, Inches(0.6), Inches(1.4), color, text=icon, font_size=22, bold=True)
        add_text(sl, x + Inches(0.7), y + Inches(0.05), Inches(5.6), Inches(0.45),
                 title, size=15, color=DARK_TEXT, bold=True)
        add_text(sl, x + Inches(0.7), y + Inches(0.5), Inches(5.6), Inches(0.9),
                 detail, size=12, color=GREY_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9: Live Demo
# ─────────────────────────────────────────────────────────────────────────────
def slide_demo(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.1), BLUE)
    add_text(sl, Inches(0.5), Inches(0.18), Inches(12), Inches(0.75),
             "Live Demonstration", size=28, color=WHITE, bold=True)

    steps = [
        ("Step 1", BLUE,   "🌐  Open Web Portal",        "http://localhost:5173",
                           "Dashboard with KPI cards showing model accuracy,\nauto-route %, and ticket volume."),
        ("Step 2", ORANGE, "📤  Upload Tickets",         "Upload → sample_tickets.xlsx",
                           "Drag and drop the Excel file. Watch the\nAnt Design Steps progress bar advance."),
        ("Step 3", GREEN,  "⚙️  Agents Processing",      "Watch the API logs",
                           "Classifier → Priority → Routing → Confidence.\nEach ticket processed in seconds."),
        ("Step 4", RGBColor(0x72, 0x2E, 0xD1), "📥  Download Results",
                           "triage_results.xlsx", "AI_Category, AI_Priority, Assigned_Team,\nConfidence_Score, Routing_Status all added."),
        ("Step 5", RGBColor(0x08, 0x97, 0x9C), "⚠️  Review Queue", "Flagged tickets",
                           "Correct AI decisions inline.\nFeedback stored for retraining."),
        ("Step 6", RED,    "📧  Check Email",             "Auto-route notification",
                           "Gmail inbox receives formatted HTML email\nfor each auto-routed P1/P2 ticket."),
    ]

    for i, (step, color, action, detail, desc) in enumerate(steps):
        row = i // 3
        col = i % 3
        x = Inches(0.4) + col * Inches(4.35)
        y = Inches(1.3) + row * Inches(2.6)
        add_rect(sl, x, y, Inches(4.1), Inches(0.45), color, text=f"{step}  —  {action}",
                 font_size=13, bold=True)
        add_rect(sl, x, y + Inches(0.45), Inches(4.1), Inches(2.0),
                 RGBColor(0x06, 0x0C, 0x1A), line_color=color, line_width=Pt(1.5))
        add_text(sl, x + Inches(0.1), y + Inches(0.6), Inches(3.9), Inches(0.4),
                 detail, size=12, color=color, bold=True)
        add_text(sl, x + Inches(0.1), y + Inches(1.05), Inches(3.9), Inches(1.2),
                 desc, size=12, color=RGBColor(0xB0, 0xC4, 0xDE))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10: KPIs & Impact
# ─────────────────────────────────────────────────────────────────────────────
def slide_kpis(prs):
    sl = blank_slide(prs)
    fill_bg(sl, LIGHT_GREY)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.1), NAVY)
    add_text(sl, Inches(0.5), Inches(0.18), Inches(12), Inches(0.75),
             "KPIs & Business Impact", size=28, color=WHITE, bold=True)

    metrics = [
        (BLUE,   "≥ 70–80%",   "Model Accuracy\n(Baseline)",     "KRA Target: ≥85% after feedback loop"),
        (GREEN,  "≥ 50%",      "Time Reduction\nin Routing",      "From ~20 min manual → <2 min automated"),
        (ORANGE, "> 85%",      "Confidence\nAuto-Route Rate",     "Tickets routed without human touch"),
        (RED,    "7 Teams",    "Covered\nAutomatically",          "Extensible via Excel config file"),
    ]

    for i, (color, value, label, sub) in enumerate(metrics):
        x = Inches(0.5) + i * Inches(3.2)
        add_rect(sl, x, Inches(1.4), Inches(2.9), Inches(1.2),
                 color, text=value, font_size=38, bold=True)
        add_text(sl, x, Inches(2.65), Inches(2.9), Inches(0.6),
                 label, size=15, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_text(sl, x, Inches(3.2), Inches(2.9), Inches(0.5),
                 sub, size=12, color=GREY_TEXT, align=PP_ALIGN.CENTER)

    # Before / After comparison
    add_rect(sl, Inches(0.4), Inches(4.0), Inches(6.0), Inches(0.45),
             RED, text="❌  BEFORE — Manual Process", font_size=14, bold=True)
    for i, line in enumerate([
        "15–30 min to triage and route a batch of tickets",
        "Inconsistent categorisation — analyst-dependent",
        "No priority logic — critical issues buried",
        "Zero visibility into routing trends or workloads",
    ]):
        add_text(sl, Inches(0.5), Inches(4.55) + i * Inches(0.45), Inches(5.8), Inches(0.4),
                 f"•  {line}", size=13, color=DARK_TEXT)

    add_rect(sl, Inches(6.9), Inches(4.0), Inches(6.0), Inches(0.45),
             GREEN, text="✅  AFTER — AI Agent System", font_size=14, bold=True)
    for i, line in enumerate([
        "< 2 min for full batch — automatic with email notification",
        "Consistent AI classification — same logic every time",
        "Priority agents assign P1–P4 based on impact keywords",
        "Real-time KPI dashboard — confidence trends, team workload",
    ]):
        add_text(sl, Inches(7.0), Inches(4.55) + i * Inches(0.45), Inches(5.8), Inches(0.4),
                 f"•  {line}", size=13, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11: Next Steps
# ─────────────────────────────────────────────────────────────────────────────
def slide_next(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)

    add_rect(sl, Inches(0), Inches(0), W, Inches(1.1), BLUE)
    add_text(sl, Inches(0.5), Inches(0.18), Inches(12), Inches(0.75),
             "Next Steps — Roadmap", size=28, color=WHITE, bold=True)

    roadmap = [
        (NAVY,   "Phase 7\n(Next)",  [
            "☁️  Azure deployment via Terraform",
            "🔄  GitHub Actions CI/CD pipeline",
            "📊  Grafana monitoring dashboards",
        ]),
        (BLUE,   "Sprint 2",  [
            "🔍  Semantic similarity (ChromaDB)",
            "📈  Improve accuracy to ≥85%",
            "🏷️  Confidence calibration tuning",
        ]),
        (GREEN,  "Sprint 3", [
            "🔗  ServiceNow / Jira integration",
            "📩  Multi-channel notifications",
            "🧪  A/B testing for routing rules",
        ]),
        (ORANGE, "Sprint 4", [
            "🤖  Self-learning retraining loop",
            "👥  Multi-tenant team support",
            "🌍  Multi-language ticket support",
        ]),
    ]

    for i, (color, phase, items) in enumerate(roadmap):
        x = Inches(0.4) + i * Inches(3.2)
        add_rect(sl, x, Inches(1.3), Inches(3.0), Inches(0.65),
                 color, text=phase, font_size=16, bold=True)
        add_rect(sl, x, Inches(1.95), Inches(3.0), Inches(4.5),
                 LIGHT_GREY, line_color=color, line_width=Pt(2))
        for j, item in enumerate(items):
            add_text(sl, x + Inches(0.15), Inches(2.1) + j * Inches(0.65),
                     Inches(2.8), Inches(0.6), item, size=14, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12: Thank You
# ─────────────────────────────────────────────────────────────────────────────
def slide_thankyou(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Inches(0), Inches(0), Inches(0.4), H, BLUE)

    add_text(sl, Inches(1.0), Inches(1.5), Inches(11), Inches(1.5),
             "Thank You", size=54, color=WHITE, bold=True)

    add_text(sl, Inches(1.0), Inches(3.2), Inches(11), Inches(0.8),
             "AI Agent Ticket Triage & Auto-Routing System",
             size=22, color=RGBColor(0x8C, 0xB8, 0xFF))

    add_text(sl, Inches(1.0), Inches(4.2), Inches(11), Inches(0.6),
             "Questions & Discussion",
             size=18, color=RGBColor(0x59, 0x6B, 0x8C))

    qas = [
        "How does the system handle tickets it has never seen before?",
        "What happens when Ollama is unavailable?",
        "How do we add new teams or change routing rules?",
        "What is the accuracy after the feedback loop runs?",
    ]
    for i, q in enumerate(qas):
        add_text(sl, Inches(1.0), Inches(5.0) + i * Inches(0.38),
                 Inches(11), Inches(0.35),
                 f"  Q:  {q}", size=12,
                 color=RGBColor(0x59, 0x6B, 0x8C))


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def convert_pptx_to_ppt(pptx_path: str) -> str | None:
    """
    Convert a .pptx file to legacy .ppt using PowerPoint COM automation.
    Requires Microsoft PowerPoint to be installed.
    Returns the .ppt path on success, None if PowerPoint is not available.
    """
    import os
    import pathlib

    pptx_abs = str(pathlib.Path(pptx_path).resolve())
    ppt_abs  = pptx_abs.replace(".pptx", ".ppt")

    try:
        import comtypes.client  # Windows only
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(pptx_abs, WithWindow=False)
        # FormatID 1 = ppSaveAsPresentation (legacy binary .ppt)
        deck.SaveAs(ppt_abs, 1)
        deck.Close()
        powerpoint.Quit()
        return ppt_abs
    except Exception as exc:
        print(f"  ⚠  COM conversion failed ({exc})")
        print("     Tip: Make sure Microsoft PowerPoint is installed, then re-run.")
        return None


def main():
    prs = new_prs()

    print("Building slides...")
    slide_title(prs)      ; print("  ✓ Slide 1: Title")
    slide_problem(prs)    ; print("  ✓ Slide 2: Problem")
    slide_solution(prs)   ; print("  ✓ Slide 3: Solution")
    slide_workflow(prs)   ; print("  ✓ Slide 4: Workflow")
    slide_tech(prs)       ; print("  ✓ Slide 5: Tech Stack")
    slide_agents(prs)     ; print("  ✓ Slide 6: Agents")
    slide_teams(prs)      ; print("  ✓ Slide 7: Teams")
    slide_kra(prs)        ; print("  ✓ Slide 8: KRA")
    slide_demo(prs)       ; print("  ✓ Slide 9: Demo")
    slide_kpis(prs)       ; print("  ✓ Slide 10: KPIs")
    slide_next(prs)       ; print("  ✓ Slide 11: Next Steps")
    slide_thankyou(prs)   ; print("  ✓ Slide 12: Thank You")

    import pathlib
    out_dir  = pathlib.Path(__file__).resolve().parent.parent
    pptx_out = str(out_dir / "AI_Ticket_Triage_Presentation.pptx")
    prs.save(pptx_out)
    print(f"\n✅  Saved → {pptx_out}  (12 slides)")

    # ── Convert to legacy .ppt via PowerPoint COM ──────────────────────────
    print("\nAttempting .ppt conversion via PowerPoint COM (skipped if unavailable)...")
    ppt_out = convert_pptx_to_ppt(pptx_out)
    if ppt_out:
        print(f"✅  Saved → {ppt_out}")
        print("\n   Both files are ready:")
        print(f"   • {pptx_out}  — modern format (recommended)")
        print(f"   • {ppt_out}   — legacy .ppt format")
    else:
        print(f"\n   The .pptx file is fully compatible with PowerPoint.")
        print(f"   To save as .ppt manually:")
        print(f"     1. Open {pptx_out} in PowerPoint")
        print(f"     2. File → Save As → 'PowerPoint 97-2003 Presentation (*.ppt)'")


if __name__ == "__main__":
    main()
